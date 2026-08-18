#!/usr/bin/env python3
"""
build_load_package.py — Generate native files + Relativity load file from mock data

Reads documents.csv for a given tier, creates actual native files (.eml, .docx,
.xlsx, .pptx, .pdf, .rsmf, etc.), and writes a Relativity .dat load file ready
for workspace import. Email and document body text is sourced from the OIDA
public S3 archive (real opioid industry OCR content). Scripted hot documents
(HOT- prefix) get hand-crafted MDL 2804 narrative content.

Output structure:
  load-packages/{tier}/
    natives/           — actual native files
    load-file.dat      — Relativity Concordance load file
    load-file.opt      — image placeholder (native-only import)
    IMPORT_README.txt  — step-by-step Relativity import instructions

Usage:
  python scripts/build_load_package.py --tier small
  python scripts/build_load_package.py --tier medium --out ./my-package/
  python scripts/build_load_package.py --tier small --no-oida   # synthetic only
  python scripts/build_load_package.py --tier small --limit 200  # first N docs only

Requirements:
  pip install python-docx openpyxl python-pptx fpdf2
"""

import argparse, csv, email.mime.multipart, email.mime.text, email.utils
import gzip, json, os, random, re, shutil, sys, time, urllib.request
from datetime import datetime
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import error_natives

try:
    from docx import Document as DocxDocument
    from docx.shared import Pt
    import openpyxl
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    from fpdf import FPDF
except ImportError:
    sys.exit("Missing dependencies. Run: pip install python-docx openpyxl python-pptx fpdf2")

DEFAULT_SEED  = 42
BUCKET_BASE   = "https://opioid-industry-documents-archive-dataset-bucket.s3.amazonaws.com"
MANIFEST_PATH = "manifest.tsv.gz"
OCR_CACHE_PATH = ".ocr-content-cache.json"
OCR_SAMPLES_PER_TYPE = 80  # number of OCR files to fetch per document category

# Relativity .dat Concordance delimiters
DAT_FIELD_SEP = chr(254)  # þ
DAT_QUOTE     = chr(255)  # ÿ
DAT_NEWLINE   = chr(174)  # ® (replaces newlines within field values)

# ── Scripted hot document content ────────────────────────────────────────

HOT_CONTENT = {
    "HOT-0000001": {
        "type": "email",
        "body": (
            "Greg,\n\n"
            "I've reviewed the Cardinal Health flag on the Cincinnati cluster. Release the order.\n\n"
            "Document it as within acceptable SOM parameters — the pharmacy has been a solid account "
            "for three years and the volume spike is attributable to Q4 stocking. Business justification "
            "is sufficient under current policy.\n\n"
            "We cannot let compliance flags slow down Q4 close. Make sure Nash has the override memo "
            "filed before end of day.\n\n"
            "Robert"
        ),
    },
    "HOT-0000002": {
        "type": "pptx",
        "slides": [
            ("Mallinckrodt Generic Opioid\nGrowth Acceleration Strategy", "McKinsey & Company — Confidential\nSeptember 2011"),
            ("Situation", "Generic oxycodone market share has plateaued at 19.4%.\nMallinckrodt has significant untapped territory in high-volume prescriber segments.\nMcKinsey analysis identifies $340M incremental revenue opportunity over 36 months."),
            ("Recommendation: Turbocharge Sales Engine", "1. Restructure territory incentive compensation to weight unit volume over market share\n2. Double down on highest-decile prescribers — top 20% generate 80% of volume\n3. Expand speaker bureau to accelerate KOL adoption in underperforming regions\n4. Align SOM thresholds to avoid unnecessary friction on growth accounts"),
            ("Territory Restructuring", "Current structure: 48 territories, quota based on market share %\nProposed: 48 territories, quota based on unit volume (oxycodone TRx)\nProjected impact: +23% unit volume in Year 1, +$118M revenue"),
            ("Implementation Timeline", "Q4 2011: Incentive compensation redesign\nQ1 2012: Territory restructuring rollout\nQ2 2012: Speaker bureau expansion — 40 new KOLs\nQ3 2012: First measurement checkpoint"),
            ("Next Steps", "Mallinckrodt leadership approval of growth framework\nMcKinsey to develop detailed implementation plan\nHR and Sales Ops alignment on new compensation structure"),
        ],
    },
    "HOT-0000003": {
        "type": "email",
        "body": (
            "Jim,\n\n"
            "FYI — see below from the DEA Diversion office. They're asking to meet on our SOM program.\n\n"
            "I've spoken with Rick at outside counsel. His recommendation is that we treat this as "
            "a routine inquiry and do not volunteer information about the Q3 2013 exception memos.\n\n"
            "I want to formally note that our SOM threshold methodology for the Cincinnati cluster "
            "and the Butner pharmacy group remains, in my view, legally exposed. I've said this before "
            "and I'm saying it again in writing.\n\n"
            "Diana\n\n"
            "---------- Forwarded message ----------\n"
            "From: DEA Diversion Control Division <diversion.inquiry@dea.gov>\n"
            "Subject: Meeting Request — Mallinckrodt SOM Program Review\n\n"
            "Mallinckrodt is requested to make available its Suspicious Order Monitoring program "
            "documentation and responsible personnel for a review meeting at the Cincinnati DEA "
            "field office. Please respond within 30 days to schedule."
        ),
    },
    "HOT-0000004": {
        "type": "email",
        "body": (
            "Natalie,\n\n"
            "Approve Foster for Q3 at $3,500/event. He delivered 6 programs last quarter with "
            "strong attendance. His Subsys volume is up 34% YoY — that's the correlation we want "
            "to see from speaker faculty.\n\n"
            "Also add Santos and Ellis to the Q4 roster at the standard rate. Loop in the IRC on "
            "prior auth support for their patient population — make sure the team knows these are "
            "priority prescribers.\n\n"
            "Alec"
        ),
    },
    "HOT-0000005": {
        "type": "docx",
        "title": "Insys Reimbursement Center — Prior Authorization Call Guide v4",
        "body": (
            "INTERNAL USE ONLY — NOT FOR DISTRIBUTION\n\n"
            "PURPOSE\n"
            "This guide provides IRC representatives with language for prior authorization calls "
            "to payer benefit lines for Subsys (fentanyl sublingual spray).\n\n"
            "CALL OPENING\n"
            "\"Hi, I'm calling from [prescriber name]'s office regarding a prior authorization "
            "request for Subsys for a patient with breakthrough cancer pain.\"\n\n"
            "WHEN ASKED ABOUT DIAGNOSIS\n"
            "Emphasize that the patient has a primary diagnosis requiring around-the-clock opioid "
            "therapy AND a separate breakthrough pain episode. If pressed on cancer diagnosis, "
            "note that the patient's condition involves significant chronic pain management needs "
            "consistent with TIRF REMS Program criteria.\n\n"
            "NOTE: Do not specify whether patient has active malignancy unless directly confirmed "
            "with prescriber. Focus on the functional description of breakthrough pain episodes.\n\n"
            "WHEN ASKED ABOUT PRIOR OPIOID THERAPY\n"
            "Confirm the patient is currently on around-the-clock opioid therapy. Reference the "
            "prescriber's clinical assessment. Do not speculate on dosage details.\n\n"
            "ESCALATION\n"
            "If payer declines, escalate to peer-to-peer review. Contact your regional Reimbursement "
            "Manager for peer-to-peer scheduling support.\n\n"
            "DOCUMENTATION\n"
            "Log all call outcomes in the IRC tracking system within 24 hours."
        ),
    },
    "HOT-0000006": {
        "type": "email",
        "body": (
            "Rick,\n\n"
            "I'm forwarding this to you rather than HR because I believe it requires legal assessment "
            "before any internal action.\n\n"
            "I received an anonymous communication through our ethics hotline raising concerns about "
            "the speaker bureau payment practices at a company we share distribution relationships "
            "with. The specific allegations involve payments to physicians that may have been "
            "structured to incentivize prescribing rather than compensate for legitimate educational "
            "services.\n\n"
            "Given the parallel nature of our own speaker program and the SOM matters currently under "
            "DEA review, I think we need privileged legal guidance on our exposure before this goes "
            "anywhere else.\n\n"
            "Please treat this as attorney-client privileged.\n\n"
            "Thomas"
        ),
    },
    "HOT-0000007": {
        "type": "email",
        "body": (
            "TO: All Mallinckrodt Employees\n"
            "FROM: Thomas Bradley, Chief Compliance Officer\n"
            "SUBJECT: LEGAL HOLD NOTICE — Opioid Litigation Matter — Immediate Action Required\n\n"
            "Mallinckrodt has received legal process in connection with litigation and regulatory "
            "matters related to opioid products. Effective immediately, you are required to preserve "
            "all documents, communications, and data that may be relevant to these matters.\n\n"
            "YOU MUST NOT DELETE, DESTROY, MODIFY, OR OVERWRITE any potentially relevant materials, "
            "including emails, attachments, reports, spreadsheets, databases, voicemails, text messages, "
            "instant messages, or any other electronic or paper records.\n\n"
            "Relevant subject matter includes:\n"
            "— Sales and marketing of opioid products (2010 to present)\n"
            "— Suspicious Order Monitoring (SOM) records and communications\n"
            "— Communications with DEA, FDA, or state regulatory authorities\n"
            "— Physician speaker bureau programs and payments\n"
            "— Distribution channel management and customer order records\n\n"
            "Acknowledge receipt of this notice by replying to this email. Failure to preserve "
            "relevant materials may result in serious legal consequences for you and the company.\n\n"
            "Questions: contact Legal at ext. 4400 or compliance@mallinckrodt.com"
        ),
    },
    "HOT-0000008": {
        "type": "xlsx",
        "title": "SOM Flag Archive Cleanup — September 2015",
        "sheets": {
            "Cleanup Log": [
                ["Record ID", "Flag Date", "Account", "Flag Type", "Volume", "Original Disposition", "Archive Action", "Archived By", "Archive Date"],
                ["SOM-2013-0847", "2013-11-14", "Cardinal Health Cincinnati", "Volume Spike", "48,000 units", "Released - BJ Memo filed", "Moved to archive", "G.Nash", "2015-09-22"],
                ["SOM-2013-0891", "2013-11-21", "Cardinal Health Butner NC", "Volume Spike", "62,000 units", "Released - BJ Memo filed", "Moved to archive", "G.Nash", "2015-09-22"],
                ["SOM-2013-0934", "2013-12-08", "McKesson Louisville", "Pattern Anomaly", "35,000 units", "Released - Regional Dir approved", "Moved to archive", "G.Nash", "2015-09-22"],
                ["SOM-2014-0012", "2014-01-17", "AmerisourceBergen Detroit", "Volume Spike", "41,000 units", "Released - BJ Memo filed", "Moved to archive", "G.Nash", "2015-09-22"],
                ["SOM-2014-0089", "2014-03-04", "Cardinal Health Cincinnati", "Repeat Flag", "55,000 units", "Released - VP Sales approved", "Moved to archive", "G.Nash", "2015-09-22"],
            ],
            "Summary": [
                ["Total records archived", "5"],
                ["Archive date", "2015-09-22"],
                ["Archived by", "G. Nash, Director SOM Compliance"],
                ["Reason for archive", "Routine records management — pre-2015 exception files"],
                ["Legal hold status at time of archive", "Not checked"],
            ],
        },
    },
    "HOT-0000009": {
        "type": "docx",
        "title": "DRAFT — Response to Ohio Attorney General Subpoena\nMallinckrodt Opioid Sales Practices Investigation",
        "body": (
            "Pursuant to the Subpoena issued by the Ohio Attorney General dated December 15, 2015, "
            "Mallinckrodt hereby responds as follows:\n\n"
            "SUBPOENA ITEM 7: Produce all records relating to Suspicious Order Monitoring (SOM) "
            "program, including all orders flagged, released, or reported to the DEA, from 2010 "
            "to present.\n\n"
            "RESPONSE TO ITEM 7:\n"
            "Mallinckrodt maintains a robust Suspicious Order Monitoring program in compliance with "
            "21 C.F.R. § 1301.74(b). Mallinckrodt will produce SOM records from January 1, 2013 "
            "to the present, subject to applicable privilege protections.\n\n"
            "[REDLINED - PRIOR DRAFT: Mallinckrodt will produce all SOM records from January 1, "
            "2010 to present — REMOVED per counsel instruction 1/12/2016]\n\n"
            "Records prior to January 1, 2013 have been archived in the ordinary course of business "
            "and retrieval would impose undue burden disproportionate to the relevance of such records.\n\n"
            "[NOTE FROM OUTSIDE COUNSEL: Do not reference the September 2015 archive cleanup in "
            "this response or in any communications with the AG's office — RG 1/13/2016]"
        ),
    },
    "HOT-0000010": {
        "type": "email",
        "body": (
            "PRIVILEGED AND CONFIDENTIAL — ATTORNEY CLIENT COMMUNICATION\n\n"
            "Rick,\n\n"
            "I just saw the news. Federal RICO charges. John Kapoor is named personally.\n\n"
            "I need to understand our exposure immediately. The speaker bureau program at Insys "
            "and the IRC prior auth practices we discussed — I need to know whether the government's "
            "theory of liability could extend to individuals who directed those programs.\n\n"
            "Also: what is the status of the DOJ civil investigation notice we received in March? "
            "And the UHC fraud unit inquiry? Are those now coordinated with the criminal case?\n\n"
            "Do not put anything in writing that isn't under privilege. Call me.\n\n"
            "Alec"
        ),
    },
    "HOT-0000011": {
        "type": "pdf",
        "body": (
            "PRIVILEGED AND CONFIDENTIAL — WORK PRODUCT\n\n"
            "MCKINSEY & COMPANY — INTERNAL MEMORANDUM\n"
            "RE: Opioid Engagement Liability — Settlement Framework Discussion\n"
            "DATE: August 30, 2017\n\n"
            "BACKGROUND\n"
            "McKinsey provided strategic consulting services to multiple pharmaceutical clients "
            "in connection with opioid products between 2004 and 2019. State attorneys general "
            "and plaintiffs' counsel in MDL 2804 have subpoenaed McKinsey work product and "
            "sought testimony from engagement personnel.\n\n"
            "LIABILITY ASSESSMENT\n"
            "Outside counsel estimates potential exposure in the range of $400M–$800M based on "
            "current settlement discussions in the MDL and parallel state AG proceedings. "
            "The principal risk is the characterization of McKinsey's opioid growth strategy "
            "recommendations as tortious conduct contributing to the epidemic.\n\n"
            "The 'turbocharge' language in the Mallinckrodt engagement deliverable (2011) and "
            "the 'Evolve to Excellence' program materials for Purdue Pharma have been cited "
            "in multiple complaints and are expected to be central to any trial.\n\n"
            "SETTLEMENT FRAMEWORK\n"
            "Outside counsel recommends proactive settlement engagement with the state AG "
            "coordinating committee. A global resolution in the $500M–$650M range is considered "
            "achievable and preferable to litigation exposure.\n\n"
            "NEXT STEPS\n"
            "Board notification required before settlement authority can be granted. "
            "Recommend scheduling board session within 60 days."
        ),
    },
    "HOT-0000012": {
        "type": "email",
        "body": (
            "Robert,\n\n"
            "Flagging an urgent issue that needs your attention today.\n\n"
            "Cardinal Health's compliance team contacted us this morning regarding an order anomaly "
            "in our Butner, NC pharmacy cluster. Three pharmacies in a 12-mile radius placed orders "
            "totaling 62,000 units of oxycodone 30mg in a 6-week period — approximately 4x their "
            "18-month average.\n\n"
            "Per our SOM procedures, this meets the threshold for mandatory reporting review. "
            "I'm escalating to you because the last two flags on this cluster were released with "
            "business justification memos, and I want senior visibility before we process this one.\n\n"
            "The accounts are Cardinal Health accounts — they're solid relationships and I don't "
            "want to damage them unnecessarily. But I also don't want us in a DEA enforcement "
            "situation because we looked the other way three times on the same cluster.\n\n"
            "What do you want me to do?\n\n"
            "Sandra"
        ),
    },
    "HOT-0000013": {
        "type": "xlsx",
        "title": "FY2013 Sales Incentive Compensation Plan — Oxycodone Territory Quotas",
        "sheets": {
            "FY2013 ICP Plan": [
                ["Territory", "Rep Name", "Base Salary", "Target TRx Quota (Oxycodone Units)", "Threshold (80%)", "Target (100%)", "Excellence (120%)", "Max Payout"],
                ["Northeast-1", "D. Morrison", "$85,000", "48,000", "$8,500", "$17,000", "$25,500", "$34,000"],
                ["Northeast-2", "R. Chen", "$82,000", "44,000", "$8,200", "$16,400", "$24,600", "$32,800"],
                ["Southeast-1", "P. Williams", "$84,000", "52,000", "$8,400", "$16,800", "$25,200", "$33,600"],
                ["Midwest-1", "K. O'Brien", "$83,000", "46,000", "$8,300", "$16,600", "$24,900", "$33,200"],
                ["Southwest-1", "J. Alvarez", "$81,000", "43,000", "$8,100", "$16,200", "$24,300", "$32,400"],
            ],
            "Notes": [
                ["FY2013 ICP — Key Changes from FY2012"],
                ["1. Quota basis changed from market share % to absolute TRx unit volume"],
                ["2. No compliance carve-out — full payout available regardless of SOM flags on territory accounts"],
                ["3. Excellence tier raised from 115% to 120% of quota"],
                ["4. McKinsey recommendation: focus incentive on volume, not market dynamics"],
                ["Approved by: Patricia Morrison, VP Marketing — January 15, 2013"],
            ],
        },
    },
}

# ── OCR content fetcher ───────────────────────────────────────────────────

def load_or_fetch_ocr_cache(use_oida, manifest_path):
    """Return a dict mapping doc type category to list of OCR text strings."""
    if not use_oida:
        return {}

    if os.path.exists(OCR_CACHE_PATH):
        print(f"  Loading OCR cache from {OCR_CACHE_PATH}...")
        with open(OCR_CACHE_PATH) as f:
            return json.load(f)

    if not os.path.exists(manifest_path):
        print(f"  WARNING: manifest not found at {manifest_path} — run scripts/fetch_manifest.py first.")
        print("  Falling back to synthetic content.")
        return {}

    print(f"  Sampling OIDA OCR files from manifest (this takes a few minutes)...")
    ocr_keys = []
    with gzip.open(manifest_path, "rt") as f:
        next(f)  # skip header
        for line in f:
            parts = line.strip().split("\t")
            if len(parts) >= 1 and parts[0].endswith(".ocr"):
                ocr_keys.append(parts[0])

    random.shuffle(ocr_keys)
    sample = ocr_keys[:OCR_SAMPLES_PER_TYPE * 6]  # fetch enough to categorize

    cache = {"email": [], "office": [], "spreadsheet": [], "presentation": [], "pdf": [], "other": []}
    fetched = 0
    for key in sample:
        if all(len(v) >= OCR_SAMPLES_PER_TYPE for v in cache.values()):
            break
        try:
            req = urllib.request.Request(f"{BUCKET_BASE}/{key}", headers={"User-Agent": "oida-registry"})
            with urllib.request.urlopen(req, timeout=10) as r:
                text = r.read(2000).decode("utf-8", errors="replace").strip()
            if not text or len(text) < 50:
                continue
            first = text.split("\n")[0].lower()
            if any(x in first for x in ["from:", "subject:", "message", "email"]):
                if len(cache["email"]) < OCR_SAMPLES_PER_TYPE: cache["email"].append(text)
            elif "sheet1" in first or "sheet " in first[:20]:
                if len(cache["spreadsheet"]) < OCR_SAMPLES_PER_TYPE: cache["spreadsheet"].append(text)
            elif any(x in text[:200].lower() for x in ["slide", "agenda", "presentation", "learning plan"]):
                if len(cache["presentation"]) < OCR_SAMPLES_PER_TYPE: cache["presentation"].append(text)
            elif any(x in first for x in ["dear ", "sincerely", "to whom", "re:", "memorandum"]):
                if len(cache["office"]) < OCR_SAMPLES_PER_TYPE: cache["office"].append(text)
            elif any(x in text[:100].lower() for x in ["page", "abstract", "introduction", "summary"]):
                if len(cache["pdf"]) < OCR_SAMPLES_PER_TYPE: cache["pdf"].append(text)
            else:
                if len(cache["other"]) < OCR_SAMPLES_PER_TYPE: cache["other"].append(text)
            fetched += 1
            if fetched % 20 == 0:
                print(f"    {fetched} OCR files fetched...")
        except Exception:
            continue

    with open(OCR_CACHE_PATH, "w") as f:
        json.dump(cache, f)
    print(f"  OCR cache built: {sum(len(v) for v in cache.values())} samples across {len(cache)} categories")
    return cache


def get_ocr_content(doc, cache):
    """Return an OCR body string for this document type."""
    ft = doc.get("File Type Category", "")
    if "Email" in ft or "Calendar" in ft:
        pool = cache.get("email", [])
    elif "Excel" in ft or "Spreadsheet" in ft:
        pool = cache.get("spreadsheet", [])
    elif "PowerPoint" in ft or "Presentation" in ft:
        pool = cache.get("presentation", [])
    elif "Word" in ft or "Document" in ft:
        pool = cache.get("office", [])
    elif "PDF" in ft:
        pool = cache.get("pdf", [])
    else:
        pool = cache.get("other", [])
    return random.choice(pool) if pool else f"[Document content — {ft}]\n\n{doc.get('Title','') or doc.get('Email Subject','')}"

# ── Native file generators ────────────────────────────────────────────────

def make_eml(doc, body):
    msg = email.mime.multipart.MIMEMultipart()
    msg["From"]       = email.utils.formataddr((doc.get("Email From",""), doc.get("Email From SMTP","")))
    msg["To"]         = email.utils.formataddr((doc.get("Email To",""),   doc.get("Email To SMTP","")))
    msg["CC"]         = doc.get("Email CC","")
    msg["Subject"]    = doc.get("Email Subject","")
    msg["Message-ID"] = doc.get("Message ID", f"<{doc['Control Number']}@mallinckrodt.com>")
    msg["Date"]       = doc.get("Date Sent","") or doc.get("Primary Date","")
    if doc.get("In Reply To"):
        msg["In-Reply-To"] = doc["In Reply To"]
    msg.attach(email.mime.text.MIMEText(body, "plain"))
    return msg.as_string()


def make_docx(doc, body):
    d = DocxDocument()
    props = d.core_properties
    props.author   = doc.get("Author","")
    props.company  = doc.get("Company","")
    props.title    = doc.get("Title","")
    if doc.get("Date Created","").strip():
        try: props.created = datetime.strptime(doc["Date Created"][:10], "%Y-%m-%d")
        except: pass
    if doc.get("Title"):
        d.add_heading(doc["Title"], 0)
    for para in body.split("\n\n"):
        d.add_paragraph(para.strip())
    import io
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def make_xlsx(doc, hot_content=None):
    import io
    wb = openpyxl.Workbook()
    if hot_content and hot_content.get("sheets"):
        # Scripted hot doc with specific sheet content
        first = True
        for sheet_name, rows in hot_content["sheets"].items():
            ws = wb.active if first else wb.create_sheet(sheet_name)
            if first: ws.title = sheet_name; first = False
            for row in rows:
                ws.append(row)
    else:
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["Field", "Value"])
        for field in ["Control Number","Custodian","Primary Date","File Type Category","Title"]:
            val = doc.get(field,"")
            if val: ws.append([field, val])
        if doc.get("Sheet Names"):
            for sname in doc["Sheet Names"].split(";")[1:]:
                wb.create_sheet(sname.strip())
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx(doc, hot_content=None):
    import io
    prs = Presentation()
    blank_layout  = prs.slide_layouts[6]
    title_layout  = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    if hot_content and hot_content.get("slides"):
        for i, (title_text, body_text) in enumerate(hot_content["slides"]):
            layout = title_layout if i == 0 else content_layout
            slide  = prs.slides.add_slide(layout)
            if slide.shapes.title: slide.shapes.title.text = title_text
            if len(slide.placeholders) > 1: slide.placeholders[1].text = body_text
    else:
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title: slide.shapes.title.text = doc.get("Title","Presentation")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Custodian: {doc.get('Custodian','')}\nDate: {doc.get('Primary Date','')[:10]}"
        for _ in range(min(int(doc.get("Slide Count",2) or 2) - 1, 4)):
            prs.slides.add_slide(blank_layout)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_pdf(doc, body):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 12)
    if doc.get("Title"):
        pdf.multi_cell(0, 6, doc["Title"][:120])
        pdf.ln(4)
    pdf.set_font("Helvetica", size=9)
    from fpdf.enums import XPos, YPos
    if doc.get("PDF Author"):
        pdf.cell(0, 5, f"Author: {doc['PDF Author']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    if doc.get("Primary Date","")[:10]:
        pdf.cell(0, 5, f"Date: {doc['Primary Date'][:10]}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(4)
    pdf.set_font("Helvetica", size=10)
    # Clean and write body
    safe_body = body.encode("latin-1", errors="replace").decode("latin-1")
    pdf.multi_cell(0, 5, safe_body[:3000])
    return pdf.output()


def make_rsmf(doc):
    """Generate a minimal RSMF JSON file."""
    data = {
        "Version": "1.0",
        "Application": doc.get("Rsmf/Application",""),
        "EventCollectionId": doc.get("Rsmf/EventCollectionId",""),
        "Participants": doc.get("Rsmf/Participants","").split("; "),
        "BeginDate": doc.get("Rsmf/BeginDate",""),
        "EndDate": doc.get("Rsmf/EndDate",""),
        "MessageCount": doc.get("Rsmf/MessageCount",0),
        "HasPlaceholders": doc.get("Rsmf/HasPlaceholders","No") == "Yes",
        "Messages": [
            {
                "Sender": doc.get("Custodian",""),
                "Timestamp": doc.get("Primary Date",""),
                "Content": f"[{doc.get('Rsmf/Application','')} conversation — {doc.get('Custodian','')}]",
                "Channel": doc.get("Conversation Topic",""),
            }
        ],
    }
    return json.dumps(data, indent=2)


def make_txt(doc, body):
    return f"Control Number: {doc.get('Control Number','')}\nDate: {doc.get('Primary Date','')[:10]}\n\n{body}"

# ── Custodian folder mapping ──────────────────────────────────────────────

_FOLDER_SAFE      = re.compile(r"[^A-Za-z0-9._-]+")
UNASSIGNED_FOLDER = "_Unassigned"


def custodian_slug(name):
    """Michael Brennan -> Michael_Brennan. Blank -> _Unassigned."""
    slug = _FOLDER_SAFE.sub("_", (name or "").strip().replace(" ", "_")).strip("_")
    return slug or UNASSIGNED_FOLDER


def native_subdir(doc):
    r"""Native directory for a document, mirroring its Processing Folder Path.

    \\Collection\Michael_Brennan\2014\01  ->  Michael_Brennan/2014/01

    The metadata column is the contract (RULES.md Rule 12): whatever it says is what
    gets built on disk, so the CSV and the package cannot drift apart. A document with
    no usable path falls back to its custodian folder; one with no custodian at all
    lands in _Unassigned.
    """
    raw   = (doc.get("Processing Folder Path") or "").replace("\\", "/")
    parts = [p for p in raw.split("/") if p and p not in (".", "..")]
    if parts and parts[0].lower() == "collection":
        parts = parts[1:]
    parts = [_FOLDER_SAFE.sub("_", p) for p in parts]
    if not parts:
        parts = [custodian_slug(doc.get("Custodian", ""))]
    return os.path.join(*parts)


def dat_native_path(abs_path, out_dir):
    """Package-relative path in the backslash form a Relativity load file expects."""
    return os.path.relpath(abs_path, out_dir).replace(os.sep, "\\")


OVERSIZED_WORD_THRESHOLD = 250_000   # past any current model context


def expand_to_word_count(body, doc):
    """Grow `body` to the Word Count the metadata claims, when that is oversized.

    Left alone for every ordinary document: only rows deliberately marked oversized
    by the edge cases cross the threshold.
    """
    try:
        target = int(str(doc.get("Word Count", "") or "0").strip())
    except ValueError:
        return body
    if target < OVERSIZED_WORD_THRESHOLD:
        return body

    seed_words = (body or "opioid distribution compliance review").split()
    if len(seed_words) < 50:
        seed_words = (seed_words * 50)[:200]

    out, n = [], 0
    para = 0
    while n < target:
        chunk = seed_words[: min(len(seed_words), target - n)]
        out.append(" ".join(chunk))
        n += len(chunk)
        para += 1
        if para % 12 == 0:
            out.append("\n\n")
    return " ".join(out)


# ── Native file dispatcher ────────────────────────────────────────────────

def generate_native(doc, cache, out_dir, flat=False, with_errors=False, error_rows=None):
    ctrl = doc["Control Number"]
    ft   = doc.get("File Type Category","")
    ext  = doc.get("File Extension","txt")

    hot  = HOT_CONTENT.get(ctrl)
    body = ""
    if hot:
        body = hot.get("body","")
    else:
        body = get_ocr_content(doc, cache)

    # Rule 13 oversized_text: the metadata's Word Count is the contract, so a
    # document claiming a million words gets a native that actually holds them.
    # Everything downstream that reads document text meets a real one here.
    body = expand_to_word_count(body, doc)

    error_rows = error_rows if error_rows is not None else []

    nat_dir = os.path.join(out_dir, "natives")
    if not flat:
        nat_dir = os.path.join(nat_dir, native_subdir(doc))
    os.makedirs(nat_dir, exist_ok=True)

    def dest(extension):
        return os.path.join(nat_dir, f"{ctrl}.{extension}")

    native_path = dest(ext)

    try:
        if ft in ("Email - MSG","Email - EML","Calendar - ICS") or (hot and hot.get("type")=="email"):
            content = make_eml(doc, body).encode("utf-8","replace")
            native_path = dest("eml")
        elif ("Word" in ft or (hot and hot.get("type")=="docx")) and "Container" not in ft:
            content = make_docx(doc, body)
            native_path = dest("docx")
        elif "Excel" in ft or "Spreadsheet" in ft or (hot and hot.get("type")=="xlsx"):
            content = make_xlsx(doc, hot)
            native_path = dest("xlsx")
        elif "PowerPoint" in ft or "Presentation" in ft or (hot and hot.get("type")=="pptx"):
            content = make_pptx(doc, hot)
            native_path = dest("pptx")
        elif "PDF" in ft or (hot and hot.get("type")=="pdf"):
            content = bytes(make_pdf(doc, body))
            native_path = dest("pdf")
        elif "RSMF" in ft or "Bloomberg" in ft:
            content = make_rsmf(doc).encode("utf-8")
            native_path = dest("rsmf")
        elif "Container" in ft or doc.get("Has Natives","") == "No":
            return None  # containers don't have natives
        else:
            content = make_txt(doc, body).encode("utf-8","replace")
            native_path = dest("txt")

        blob = content if isinstance(content, bytes) else content.encode("utf-8","replace")
        with open(native_path, "wb") as f:
            f.write(blob)

        # Rule 12: a document the metadata flags as an error gets a native that
        # actually fails that way. Fabricated in place, over the healthy file.
        record = error_natives.fabricate(doc, native_path, blob) if with_errors else None
        if record is not None:
            record["Native File"] = dat_native_path(native_path, out_dir)
            error_rows.append(record)

        return dat_native_path(native_path, out_dir)

    except Exception as e:
        # Fallback: write a txt placeholder. Never rescue a deliberately broken
        # file — if fabrication failed, that is a bug worth seeing.
        if with_errors and error_natives.scenario_for(doc) is not None:
            raise
        fallback = dest("txt")
        with open(fallback, "w", encoding="utf-8") as f:
            f.write(f"[{ft}]\n\n{body[:500]}")
        return dat_native_path(fallback, out_dir)


# ── Custodian data source sheet ───────────────────────────────────────────

CUSTODIAN_SOURCE_COLUMNS = [
    "Custodian","Custodian Email","Custodian Org","Custodian Department",
    "Data Source Folder","Documents","Natives Written","Native Bytes",
]


def write_custodian_sources(stats, out_dir, flat):
    """One row per custodian: the setup sheet for Relativity processing data sources."""
    path = os.path.join(out_dir, "custodian-sources.csv")
    with open(path, "w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(CUSTODIAN_SOURCE_COLUMNS)
        for name in sorted(stats):
            st = stats[name]
            folder = "natives" if flat else "natives\\" + st["folder"]
            w.writerow([name, st["email"], st["org"], st["dept"], folder,
                        st["docs"], st["natives"], st["bytes"]])
    return path


# ── Relativity .dat load file writer ─────────────────────────────────────

# Map documents.csv column names to Relativity standard field names
RELATIVITY_FIELD_MAP = {
    "Control Number":         "BegDoc#",
    "Control Number":         "EndDoc#",     # same value for single-page docs
    "Family ID":              "BegAttach",
    "Family ID":              "EndAttach",
    "Custodian":              "Custodian",
    "Custodian Email":        "Custodian Email",
    "Custodian Org":          "Custodian Org",
    "File Name":              "File Name",
    "File Extension":         "File Type",
    "File Size (bytes)":      "File Size",
    "Primary Date":           "Date",
    "Email From":             "From",
    "Email From SMTP":        "From (SMTP)",
    "Email To":               "To",
    "Email To SMTP":          "To (SMTP)",
    "Email CC":               "CC",
    "Email Subject":          "Subject",
    "Date Sent":              "Date Sent",
    "Date Received":          "Date Received",
    "Message ID":             "Message ID",
    "Has Attachments":        "Has Attachments",
    "Attachment Count":       "Attachment Count",
    "Email Thread ID":        "Email Thread ID",
    "Email Threading Inclusive": "Email Threading Inclusive",
    "Conversation Topic":     "Conversation Topic",
    "Author":                 "Author",
    "Title":                  "Title",
    "Company":                "Company",
    "Page Count":             "Page Count",
    "Workflow Stage":         "Workflow Stage",
    "Responsiveness":         "Responsive",
    "Privilege":              "Privileged",
    "Privilege Reason":       "Privilege Reason",
    "Hot Doc":                "Hot Doc",
    "Issue Tags":             "Issue Tags",
    "Bates Begin":            "BegBates",
    "Bates End":              "EndBates",
    "Production Set":         "Production Set",
    "Redacted":               "Redacted",
    "TAR Score":              "TAR Score",
    "AL Predicted Relevant":  "AL Predicted Relevant",
    "Batch Name":             "Batch Name",
    "Batch Status":           "Batch Status",
    "Reviewer":               "Reviewer",
    "Narrative Phase":        "Narrative Phase",
    "Narrative Phase Name":   "Narrative Phase Name",
    "Dedup Method":           "Dedup Method",
    "MD5 Hash":               "MD5 Hash",
    "OCR Flag":               "OCR Flag",
    "Rsmf/Application":       "Rsmf Application",
    "Rsmf/Participants":      "Rsmf Participants",
    "Rsmf/MessageCount":      "Rsmf Message Count",
    "NativeFilePath":         "NativeFilePath",
}

DAT_COLUMNS = [
    "BegDoc#","EndDoc#","BegAttach","EndAttach","Custodian","Custodian Email",
    "Custodian Org","File Name","File Type","File Size","Date","From","From (SMTP)",
    "To","To (SMTP)","CC","Subject","Date Sent","Date Received","Message ID",
    "Has Attachments","Attachment Count","Email Thread ID","Email Threading Inclusive",
    "Conversation Topic","Author","Title","Company","Page Count",
    "Workflow Stage","Responsive","Privileged","Privilege Reason","Hot Doc","Issue Tags",
    "BegBates","EndBates","Production Set","Redacted","TAR Score","AL Predicted Relevant",
    "Batch Name","Batch Status","Reviewer","Narrative Phase","Narrative Phase Name",
    "Dedup Method","MD5 Hash","OCR Flag","Rsmf Application","Rsmf Participants",
    "Rsmf Message Count","Processing Status","Processing Error Type","NativeFilePath",
]


def dat_row(values):
    """Format a list of values as a Relativity .dat row."""
    def clean(v):
        s = str(v) if v is not None else ""
        return s.replace(DAT_QUOTE, "").replace(DAT_FIELD_SEP, " ").replace("\n", DAT_NEWLINE).replace("\r","")
    return DAT_FIELD_SEP.join(DAT_QUOTE + clean(v) + DAT_QUOTE for v in values) + "\n"


# Declarative mapping: .dat column name → (source key in doc dict, optional transform)
# None transform = direct doc.get(key, ""); callable transform receives the full doc.
_COLUMN_MAP = {
    "BegDoc#":                   ("Control Number",          None),
    "EndDoc#":                   ("Control Number",          None),
    "Custodian":                 ("Custodian",               None),
    "Custodian Email":           ("Custodian Email",         None),
    "Custodian Org":             ("Custodian Org",           None),
    "File Name":                 ("File Name",               None),
    "File Type":                 ("File Extension",          None),
    "File Size":                 ("File Size (bytes)",       None),
    "Date":                      ("Primary Date",            lambda d: d.get("Primary Date","")[:10]),
    "From":                      ("Email From",              None),
    "From (SMTP)":               ("Email From SMTP",         None),
    "To":                        ("Email To",                None),
    "To (SMTP)":                 ("Email To SMTP",           None),
    "CC":                        ("Email CC",                None),
    "Subject":                   ("Email Subject",           lambda d: d.get("Email Subject","") or d.get("Title","")),
    "Date Sent":                 ("Date Sent",               lambda d: d.get("Date Sent","")[:10] if d.get("Date Sent") else ""),
    "Date Received":             ("Date Received",           lambda d: d.get("Date Received","")[:10] if d.get("Date Received") else ""),
    "Message ID":                ("Message ID",              None),
    "Has Attachments":           ("Has Attachments",         None),
    "Attachment Count":          ("Attachment Count",        None),
    "Email Thread ID":           ("Email Thread ID",         None),
    "Email Threading Inclusive": ("Email Threading Inclusive",None),
    "Conversation Topic":        ("Conversation Topic",      None),
    "Author":                    ("Author",                  None),
    "Title":                     ("Title",                   None),
    "Company":                   ("Company",                 None),
    "Page Count":                ("Page Count",              None),
    "Workflow Stage":            ("Workflow Stage",          None),
    "Responsive":                ("Responsiveness",          None),
    "Privileged":                ("Privilege",               None),
    "Privilege Reason":          ("Privilege Reason",        None),
    "Hot Doc":                   ("Hot Doc",                 None),
    "Issue Tags":                ("Issue Tags",              None),
    "BegBates":                  ("Bates Begin",             None),
    "EndBates":                  ("Bates End",               None),
    "Production Set":            ("Production Set",          None),
    "Redacted":                  ("Redacted",                None),
    "TAR Score":                 ("TAR Score",               None),
    "AL Predicted Relevant":     ("AL Predicted Relevant",   None),
    "Batch Name":                ("Batch Name",              None),
    "Batch Status":              ("Batch Status",            None),
    "Reviewer":                  ("Reviewer",                None),
    "Narrative Phase":           ("Narrative Phase",         None),
    "Narrative Phase Name":      ("Narrative Phase Name",    None),
    "Dedup Method":              ("Dedup Method",            None),
    "MD5 Hash":                  ("MD5 Hash",                None),
    "OCR Flag":                  ("OCR Flag",                None),
    "Rsmf Application":          ("Rsmf/Application",        None),
    "Rsmf Participants":         ("Rsmf/Participants",       None),
    "Rsmf Message Count":        ("Rsmf/MessageCount",       None),
    "Processing Status":         ("Processing Status",       None),
    "Processing Error Type":     ("Processing Error Type",   None),
}


def doc_to_dat_row(doc, native_rel_path, families_by_doc, native_bytes=None):
    fam    = families_by_doc.get(doc.get("Control Number",""), {})
    values = []
    for col in DAT_COLUMNS:
        if col == "BegAttach":    v = fam.get("beg_attach","")
        elif col == "EndAttach":  v = fam.get("end_attach","")
        elif col == "NativeFilePath": v = native_rel_path or ""
        # The size on disk is the truth; the metadata describes a file that was
        # never written, and every size-based check downstream needs the real one.
        elif col == "File Size" and native_bytes is not None: v = native_bytes
        elif col in _COLUMN_MAP:
            src_key, transform = _COLUMN_MAP[col]
            v = transform(doc) if transform else doc.get(src_key,"")
        else:
            v = ""
        values.append(v)
    return values


def build_family_index(families_path):
    """Build a dict: control_number → {beg_attach, end_attach}."""
    if not os.path.exists(families_path):
        return {}
    with open(families_path) as f:
        families = json.load(f)
    index = {}
    for fam in families:
        beg = fam.get("parent_doc_id","")
        children = fam.get("children",[])
        end = children[-1] if children else beg
        for doc_id in [beg] + children:
            index[doc_id] = {"beg_attach": beg, "end_attach": end}
    return index

# ── Import readme ─────────────────────────────────────────────────────────

def custodian_readme_block(stats, flat):
    """The custodian folder listing embedded in IMPORT_README.txt."""
    if flat:
        total = sum(st["docs"] for st in stats.values())
        return ("  natives\\   all {:,} documents in one folder (--flat).\n"
                "            Custodian is NOT derivable from the folder structure here;\n"
                "            use PATH B, or rebuild without --flat.".format(total))
    lines = []
    for name in sorted(stats):
        st = stats[name]
        lines.append("  natives\\{:<20} {:>7,} docs   {:>8.1f} MB   {}".format(
            st["folder"], st["docs"], st["bytes"]/1e6, name))
    lines.append("")
    lines.append("Each custodian folder is further split by year and month, mirroring the")
    lines.append("Processing Folder Path column in documents.csv. The folder structure and")
    lines.append("that column always agree, so either can be treated as the source of truth.")
    return "\n".join(lines)


IMPORT_README = """RELATIVITY IMPORT INSTRUCTIONS
====================================

This package contains:
  natives/               native files, organised into one folder per custodian
  load-file.dat          Relativity Concordance load file (metadata + native paths)
  load-file.opt          image load file placeholder (native-only import)
  custodian-sources.csv  one row per custodian: the data source setup sheet

CUSTODIAN FOLDERS
{custodian_block}

Pick ONE of the two paths below. PATH A exercises Processing and is the right
choice for testing the raw data workflow. PATH B skips Processing and loads the
metadata as-is.


PATH A — PROCESS AS RAW DATA (recommended for workflow testing)
---------------------------------------------------------------

STEP A1 — Copy the package to your Relativity file server
  Place this entire folder on the server at a path Relativity can access.
  Example: \\\\fileserver\\LoadFiles\\MDL2804-{tier}\\

STEP A2 — Create a processing profile
  Processing → Profiles → New. Default settings are fine to start.

STEP A3 — Create a processing set with ONE DATA SOURCE PER CUSTODIAN
  Processing → Processing Sets → New, then add a data source for each row in
  custodian-sources.csv:

    Data Source Folder  →  the folder in the "Data Source Folder" column
    Custodian           →  the "Custodian" column, create the entity if needed
    Document numbering  →  your choice

  This is the whole point of the folder layout: custodian assignment comes from
  the folder structure, so you never hand-sort files or hand-map custodians.

STEP A4 — Run discovery, then publish
  Check the errors tab before publishing. In a default package every file is
  expected to process cleanly.


PATH B — IMPORT THE LOAD FILE (metadata already populated)
-----------------------------------------------------------

STEP B1 — Copy the package to a location Relativity can reach (as in A1).

STEP B2 — Workspace → Import → Relativity Load File → select load-file.dat

STEP B3 — Field mapping
  The .dat file uses Concordance delimiters:
    Column separator: þ (ASCII 254)
    Text qualifier:   ÿ (ASCII 255)
    Newline in field: ® (ASCII 174)

  Map these .dat columns to Relativity fields:
    BegDoc#              → Control Number
    Custodian            → Custodian
    Custodian Org        → Custodian Org (custom text field)
    Responsive           → Responsiveness (single choice)
    Privileged           → Privilege (single choice)
    Hot Doc              → Hot Doc (yes/no)
    Issue Tags           → Issue Tags (multi-choice or long text)
    Narrative Phase      → Narrative Phase (number)
    Narrative Phase Name → Narrative Phase Name (text)
    TAR Score            → TAR Score (decimal)
    NativeFilePath       → (mapped to native file upload)

STEP B4 — Set the native file path base
  When prompted for the native file path, set the base path to the location of
  this package on the file server. NativeFilePath holds package-relative paths
  like: natives\\Michael_Brennan\\2014\\01\\DOC-0000318.docx


VERIFYING
  After loading, run this search to find the scripted hot docs:
    Control Number StartsWith "HOT-"

  These 8–13 documents are the key evidentiary moments in the MDL 2804 story.
  See mock-data/DEMO_GUIDE.md for a full walkthrough.

  To confirm custodian assignment worked, group the document list by Custodian
  and compare the counts against custodian-sources.csv.

QUESTIONS
  See CONTRIBUTING.md or open a GitHub issue at:
  https://github.com/nickmanoogian/oida-registry
"""


def apply_error_rate(all_docs, error_rate, seed):
    """Promote extra clean documents to errors until the target rate is reached.

    Rule 6 sets the baseline distribution (~8% in every tier). A bug bash may want
    the failure paths hit far harder than production ever would, so extra documents
    are promoted using the same mix of error types already present, and the metadata
    is mutated so the load file and the natives agree.
    """
    if error_rate is None:
        return 0
    errored = [d for d in all_docs if (d.get("Processing Error Type") or "").strip()]
    target  = int(len(all_docs) * error_rate)
    if target <= len(errored):
        return 0

    types = [(d.get("Processing Error Type") or "").strip() for d in errored]
    types = [t for t in types if t and t not in error_natives.NOT_FABRICABLE] or ["Corrupt File"]

    rng   = random.Random(seed)
    clean = [d for d in all_docs
             if not (d.get("Processing Error Type") or "").strip()
             and d.get("Has Natives","") != "No"]
    rng.shuffle(clean)

    promoted = 0
    for doc in clean[: target - len(errored)]:
        doc["Processing Error Type"] = rng.choice(types)
        doc["Processing Status"]     = "Error"
        doc["Workflow Stage"]        = "Pre-Review: Processing Error"
        promoted += 1
    return promoted


def write_expected_errors(error_rows, out_dir):
    path = os.path.join(out_dir, "EXPECTED_ERRORS.csv")
    with open(path, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=error_natives.EXPECTED_ERROR_COLUMNS)
        w.writeheader()
        for row in sorted(error_rows, key=lambda r: r["Control Number"]):
            w.writerow(row)
    return path


def edge_cases_readme_block(edge_path):
    """The starved-documents section appended to IMPORT_README.txt."""
    if not edge_path:
        return ""
    with open(edge_path) as f:
        scenarios = json.load(f)["scenarios"]
    total = sum(v.get("count", 0) for v in scenarios.values())
    lines = ["", "", "DOCUMENTS THAT PROCESS CLEANLY AND ARE STILL INCOMPLETE",
             "=" * 54, "",
             f"{total} documents in this package are missing something a feature",
             "depends on. They are not processing errors: they will import and",
             "process without complaint. That is the point.", ""]
    for name in sorted(scenarios):
        v = scenarios[name]
        lines.append(f"  {v.get('count',0):>4}  {name:<22} starves {v.get('starves','')}")
    lines += ["",
              "edge-cases.json lists every one by control number.",
              "",
              "WHAT TO LOOK FOR",
              "  Are these documents counted, excluded, or silently dropped? Does a",
              "  per-custodian view acknowledge the ones with no custodian? Does a",
              "  timeline survive a 1601 date? Does a language summary survive a",
              "  document in three languages?", ""]
    return "\n".join(lines)


def expected_errors_readme_block(error_rows):
    """The intentionally-broken-files section appended to IMPORT_README.txt."""
    if not error_rows:
        return ""
    by_scenario = {}
    for r in error_rows:
        by_scenario.setdefault(r["Scenario"], []).append(r)
    lines = ["", "", "INTENTIONALLY BROKEN FILES IN THIS PACKAGE",
             "=" * 42, "",
             f"This package was built with --with-errors. {len(error_rows)} documents contain",
             "natives designed to fail processing. This is deliberate.", ""]
    for scenario in sorted(by_scenario):
        rows = by_scenario[scenario]
        flag = "" if rows[0]["Guaranteed"] == "yes" else "   (not guaranteed)"
        lines.append(f"  {len(rows):>4}  {scenario}{flag}")
    lines += ["",
              f"Encrypted files use the password: {error_natives.PACKAGE_PASSWORD}",
              "Add it to the Relativity password bank to test the recovery path,",
              "or leave it out to test the failure path.",
              "",
              "EXPECTED_ERRORS.csv lists every one: control number, custodian, native",
              "file, how it was built, the error Relativity is expected to report, and",
              "whether that outcome is guaranteed. Rows marked Guaranteed = no depend on",
              "engine or worker configuration rather than on the file.",
              "",
              "AFTER PROCESSING",
              "  1. Open the Processing Set error report.",
              "  2. Compare it against EXPECTED_ERRORS.csv.",
              "  3. Record anything that differs. Two kinds of finding matter:",
              "       - a file expected to fail that processed cleanly",
              "       - a file that failed with a different error than expected",
              "  4. Then look at what the downstream feature does with the failed set:",
              "     are those documents counted, excluded, or silently dropped?",
              "",
              "DO NOT report the processing errors themselves as bugs. They are the",
              "point. Report what the product does with them.", ""]
    return "\n".join(lines)


# ── Main ──────────────────────────────────────────────────────────────────

def build(tier_name, tier_dir, out_dir, use_oida, limit, seed, flat=False,
          with_errors=False, error_rate=None):
    random.seed(seed)

    docs_path    = os.path.join(tier_dir, "documents.csv")
    families_path = os.path.join(tier_dir, "email-families.json")

    if not os.path.exists(docs_path):
        sys.exit(f"ERROR: {docs_path} not found. Run generate_mock_metadata.py --tier {tier_name} first.")

    with open(docs_path, encoding="utf-8") as f:
        all_docs = list(csv.DictReader(f))

    if limit:
        # Always include all HOT- docs, then fill with limit from the rest
        hot_docs  = [d for d in all_docs if d["Control Number"].startswith("HOT-")]
        rest      = [d for d in all_docs if not d["Control Number"].startswith("HOT-")]
        all_docs  = hot_docs + rest[:max(0, limit - len(hot_docs))]

    print(f"\n{'='*60}\n  Building load package: MDL 2804 {tier_name.upper()}")
    print(f"  Documents:  {len(all_docs):,}")
    print(f"  Output:     {out_dir}")
    print(f"  OIDA OCR:   {'yes' if use_oida else 'no (synthetic)'}")
    print(f"  Layout:     {'flat natives/' if flat else 'natives/{custodian}/{year}/{month}'}")
    if with_errors:
        promoted = apply_error_rate(all_docs, error_rate, seed)
        flagged  = sum(1 for d in all_docs if (d.get("Processing Error Type") or "").strip())
        print(f"  Errors:     on — {flagged:,} documents flagged"
              + (f" ({promoted:,} promoted to hit {error_rate:.0%})" if promoted else ""))
    print(f"{'='*60}\n")

    Path(out_dir).mkdir(parents=True, exist_ok=True)
    Path(os.path.join(out_dir, "natives")).mkdir(exist_ok=True)

    # Load OCR content
    cache = load_or_fetch_ocr_cache(use_oida, MANIFEST_PATH)

    # Build family index for attach ranges
    families_by_doc = build_family_index(families_path)

    # Generate natives + build .dat rows
    dat_rows = []
    skipped  = 0
    natives_written = 0
    cust_stats = {}
    error_rows = []
    t0 = time.time()

    for i, doc in enumerate(all_docs):
        native_path = generate_native(doc, cache, out_dir, flat=flat,
                                      with_errors=with_errors, error_rows=error_rows)
        native_bytes = (os.path.getsize(os.path.join(out_dir, native_path.replace("\\", os.sep)))
                        if native_path else None)
        values = doc_to_dat_row(doc, native_path, families_by_doc, native_bytes)
        dat_rows.append(values)

        cust = (doc.get("Custodian") or "").strip() or "(unassigned)"
        st   = cust_stats.setdefault(cust, {
            "email":  doc.get("Custodian Email",""),
            "org":    doc.get("Custodian Org",""),
            "dept":   doc.get("Custodian Department",""),
            "folder": native_subdir(doc).split(os.sep)[0],
            "docs": 0, "natives": 0, "bytes": 0,
        })
        st["docs"] += 1
        if native_path:
            natives_written += 1
            st["natives"] += 1
            st["bytes"]   += native_bytes
        else:
            skipped += 1

        if (i+1) % 100 == 0:
            elapsed = time.time() - t0
            rate    = (i+1) / elapsed
            print(f"  {i+1:,}/{len(all_docs):,} docs  ({rate:.0f}/s)  {natives_written:,} natives written")

    # Write .dat
    dat_path = os.path.join(out_dir, "load-file.dat")
    with open(dat_path, "w", encoding="utf-8", newline="") as f:
        f.write(dat_row(DAT_COLUMNS))
        for row in dat_rows:
            f.write(dat_row(row))
    dat_mb = os.path.getsize(dat_path) / 1e6

    # Write .opt placeholder
    opt_path = os.path.join(out_dir, "load-file.opt")
    with open(opt_path, "w") as f:
        f.write("# Relativity image load file\n")
        f.write("# This package uses native-only import — no TIFF images included.\n")
        f.write("# Import using load-file.dat for native file loading.\n")

    # Write the custodian data source sheet
    write_custodian_sources(cust_stats, out_dir, flat)

    # Write the expected-errors manifest
    if with_errors:
        write_expected_errors(error_rows, out_dir)

    # Carry the edge-case manifest across from the metadata tier. Without it the
    # package has the starved documents but no map of which ones are deliberate,
    # which is the same hole EXPECTED_ERRORS.csv exists to close.
    edge_src = os.path.join(tier_dir, "edge-cases.json")
    edge_count = 0
    if os.path.exists(edge_src):
        shutil.copyfile(edge_src, os.path.join(out_dir, "edge-cases.json"))
        with open(edge_src) as f:
            edge_count = sum(v.get("count", 0) for v in json.load(f)["scenarios"].values())

    # Write import readme
    readme_path = os.path.join(out_dir, "IMPORT_README.txt")
    with open(readme_path, "w") as f:
        f.write(IMPORT_README.replace("{tier}", tier_name)
                             .replace("{custodian_block}", custodian_readme_block(cust_stats, flat))
                + expected_errors_readme_block(error_rows)
                + edge_cases_readme_block(edge_src if edge_count else None))

    print(f"\n  Done in {time.time()-t0:.0f}s")
    print(f"  Natives:    {natives_written:,} files ({skipped:,} documents have no native)")
    print(f"  load-file.dat: {dat_mb:.1f} MB ({len(dat_rows):,} rows, {len(DAT_COLUMNS)} fields)")
    if with_errors:
        guaranteed = sum(1 for r in error_rows if r["Guaranteed"] == "yes")
        if edge_count:
            print(f"  Starved:    {edge_count:,} documents missing an input -> edge-cases.json")
        print(f"  Broken:     {len(error_rows):,} natives fabricated "
              f"({guaranteed:,} guaranteed) -> EXPECTED_ERRORS.csv")
        skipped_kinds = {(d.get("Processing Error Type") or "").strip()
                         for d in all_docs} & set(error_natives.NOT_FABRICABLE)
        for kind in sorted(skipped_kinds):
            print(f"    not fabricated: {kind} — {error_natives.NOT_FABRICABLE[kind]}")
    print(f"  Custodians: {len(cust_stats)} -> custodian-sources.csv")
    for name in sorted(cust_stats):
        st = cust_stats[name]
        folder = "natives" if flat else os.path.join("natives", st["folder"])
        print(f"    {name:<22} {st['docs']:>7,} docs  {st['bytes']/1e6:>7.1f} MB  {folder}")
    print(f"  Package:    {out_dir}")


def main():
    p = argparse.ArgumentParser(description="Build Relativity native file load package from mock data")
    p.add_argument("--tier",     required=True, choices=["small","medium","large"])
    p.add_argument("--dir",      default=None,  help="Source tier directory (default: mock-data/{tier}/)")
    p.add_argument("--out",      default=None,  help="Output directory (default: load-packages/{tier}/)")
    p.add_argument("--no-oida",  action="store_true", help="Use synthetic content instead of OIDA OCR")
    p.add_argument("--limit",    type=int, default=None, help="Only process first N documents (useful for testing)")
    p.add_argument("--seed",     type=int, default=DEFAULT_SEED)
    p.add_argument("--with-errors", action="store_true",
                   help="Fabricate natives that genuinely fail processing for every "
                        "document flagged Processing Status = Error")
    p.add_argument("--error-rate", type=float, default=None,
                   help="Promote extra documents to errors until this fraction is reached "
                        "(e.g. 0.25). Requires --with-errors")
    p.add_argument("--flat",     action="store_true",
                   help="Write every native into one natives/ directory instead of "
                        "natives/{custodian}/{year}/{month}/")
    args = p.parse_args()

    tier_dir = args.dir or os.path.join("mock-data", args.tier)
    out_dir  = args.out or os.path.join("load-packages", args.tier)
    if args.error_rate is not None and not args.with_errors:
        p.error("--error-rate requires --with-errors")

    build(args.tier, tier_dir, out_dir, not args.no_oida, args.limit, args.seed, args.flat,
          args.with_errors, args.error_rate)

if __name__ == "__main__":
    main()
